"""PPTX manipulation service using python-pptx"""
import io
import os
import logging
import tempfile
from typing import Dict, List, Optional, BinaryIO
from uuid import UUID
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.infrastructure.text_layout import TextLayoutService

logger = logging.getLogger(__name__)


class PPTXService:
    """Service for creating and manipulating PPTX files"""
    
    @staticmethod
    def create_empty_presentation() -> bytes:
        """
        Create an empty PPTX presentation with one blank slide
        
        Returns:
            PPTX file as bytes
        """
        prs = PPTXPresentation()
        
        # Log initial slide count
        logger.info(f"Initial slide count after PPTXPresentation(): {len(prs.slides)}")
        
        # Check if presentation already has slides
        if len(prs.slides) == 0:
            # No default slide, add a blank one
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            prs.slides.add_slide(blank_slide_layout)
            logger.info(f"Added blank slide. Total slides: {len(prs.slides)}")
        else:
            # Already has slides, don't add more
            logger.info(f"Presentation already has {len(prs.slides)} slide(s), not adding more")
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()
    
    @staticmethod
    def load_presentation(pptx_data: bytes) -> PPTXPresentation:
        """
        Load a PPTX presentation from bytes
        
        Args:
            pptx_data: PPTX file data as bytes
            
        Returns:
            Presentation object
        """
        input_stream = io.BytesIO(pptx_data)
        return PPTXPresentation(input_stream)
    
    @staticmethod
    def save_presentation(prs: PPTXPresentation) -> bytes:
        """
        Save a PPTX presentation to bytes
        
        Args:
            prs: Presentation object
            
        Returns:
            PPTX file as bytes
        """
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()
    
    @staticmethod
    def add_slide(prs: PPTXPresentation, layout_index: int = 6) -> int:
        """
        Add a new slide to the presentation
        
        Args:
            prs: Presentation object
            layout_index: Index of slide layout (6 = blank)
            
        Returns:
            Index of the new slide
        """
        slide_layout = prs.slide_layouts[layout_index]
        prs.slides.add_slide(slide_layout)
        return len(prs.slides) - 1
    
    @staticmethod
    def copy_shapes_from_block(
        target_prs: PPTXPresentation,
        block_pptx_data: bytes,
        slide_index: int
    ) -> bool:
        """
        Copy all shapes from a block PPTX to a target presentation slide
        
        Args:
            target_prs: Target presentation
            block_pptx_data: Block PPTX file data
            slide_index: Target slide index
            
        Returns:
            True if successful
        """
        try:
            # Load block presentation
            block_prs = PPTXService.load_presentation(block_pptx_data)
            
            # Get first slide from block (blocks are single-slide templates)
            if len(block_prs.slides) == 0:
                logger.warning("Block has no slides")
                return False
            
            source_slide = block_prs.slides[0]
            target_slide = target_prs.slides[slide_index]
            
            # Copy all shapes from source to target
            for shape in source_slide.shapes:
                PPTXService._copy_shape(shape, target_slide)
            
            logger.info(f"Copied {len(source_slide.shapes)} shapes to slide {slide_index}")
            return True
            
        except Exception as e:
            logger.error(f"Error copying shapes from block: {e}")
            return False
    
    @staticmethod
    def _copy_shape(source_shape, target_slide):
        """
        Copy a single shape to target slide
        
        Note: This is a simplified implementation. 
        For production, consider using more sophisticated shape cloning.
        """
        try:
            # Handle different shape types
            if source_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                # Copy text box
                textbox = target_slide.shapes.add_textbox(
                    source_shape.left,
                    source_shape.top,
                    source_shape.width,
                    source_shape.height
                )
                if source_shape.has_text_frame:
                    textbox.text_frame.text = source_shape.text_frame.text
                    # Copy text formatting
                    for src_para, tgt_para in zip(source_shape.text_frame.paragraphs, textbox.text_frame.paragraphs):
                        tgt_para.alignment = src_para.alignment
                        for src_run, tgt_run in zip(src_para.runs, tgt_para.runs):
                            tgt_run.font.size = src_run.font.size
                            tgt_run.font.bold = src_run.font.bold
                            tgt_run.font.italic = src_run.font.italic
                            if src_run.font.color.rgb:
                                tgt_run.font.color.rgb = src_run.font.color.rgb
            
            elif source_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Copy picture
                image_stream = io.BytesIO(source_shape.image.blob)
                target_slide.shapes.add_picture(
                    image_stream,
                    source_shape.left,
                    source_shape.top,
                    source_shape.width,
                    source_shape.height
                )
            
            elif source_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Copy auto shape
                shape = target_slide.shapes.add_shape(
                    source_shape.auto_shape_type,
                    source_shape.left,
                    source_shape.top,
                    source_shape.width,
                    source_shape.height
                )
                if source_shape.has_text_frame and shape.has_text_frame:
                    shape.text_frame.text = source_shape.text_frame.text
            
            elif source_shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                # Copy table
                table = target_slide.shapes.add_table(
                    len(source_shape.table.rows),
                    len(source_shape.table.columns),
                    source_shape.left,
                    source_shape.top,
                    source_shape.width,
                    source_shape.height
                ).table
                
                # Copy cell contents
                for row_idx, row in enumerate(source_shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        target_cell = table.cell(row_idx, col_idx)
                        target_cell.text = cell.text
            
            else:
                logger.warning(f"Unsupported shape type: {source_shape.shape_type}")
                
        except Exception as e:
            logger.error(f"Error copying shape: {e}")
    
    @staticmethod
    def replace_placeholders(
        prs: PPTXPresentation,
        slide_index: int,
        replacements: Dict[str, str]
    ) -> bool:
        """
        Replace placeholders in a slide with actual values (старый метод)
        
        Args:
            prs: Presentation object
            slide_index: Index of slide to modify
            replacements: Dictionary of {placeholder: value}
            
        Returns:
            True if successful
        """
        try:
            slide = prs.slides[slide_index]
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    # Check each paragraph
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Replace placeholders
                            for placeholder, value in replacements.items():
                                placeholder_text = f"{{{{{placeholder}}}}}"
                                if placeholder_text in run.text:
                                    run.text = run.text.replace(placeholder_text, str(value))
                                    logger.info(f"Replaced {placeholder_text} with {value}")
            
            return True
            
        except Exception as e:
            logger.error(f"Error replacing placeholders: {e}")
            return False
    
    @staticmethod
    def fill_template_with_data(
        prs: PPTXPresentation,
        slide_index: int,
        fields_data: List[Dict]
    ) -> bool:
        """
        Заполнить слайд данными с автоматическим размещением
        
        Args:
            prs: Presentation object
            slide_index: Index of slide to modify
            fields_data: Список данных полей в формате:
                [{
                    "text": "Содержимое",
                    "font_metadata": {
                        "name": "Arial",
                        "size": 18,
                        "color": "#000000",
                        "bold": False,
                        "italic": False,
                        "alignment": "LEFT"
                    },
                    "order_index": 0
                }, ...]
            
        Returns:
            True if successful
        """
        try:
            slide = prs.slides[slide_index]
            
            # Логируем все shapes на слайде до удаления
            logger.info(f"Slide {slide_index} has {len(slide.shapes)} shapes before cleanup")
            for i, shape in enumerate(slide.shapes):
                shape_info = f"Shape {i}: type={shape.shape_type}"
                if hasattr(shape, 'name'):
                    shape_info += f", name={shape.name}"
                if shape.has_text_frame:
                    shape_info += f", text={shape.text_frame.text[:30] if shape.text_frame.text else 'empty'}"
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    try:
                        shape_info += f", placeholder_type={shape.placeholder_format.type}"
                    except:
                        shape_info += ", placeholder_type=unknown"
                # Проверяем fill (заливку)
                try:
                    if hasattr(shape, 'fill'):
                        fill_type = shape.fill.type
                        shape_info += f", fill_type={fill_type}"
                        if fill_type == 1:  # SOLID
                            try:
                                color = shape.fill.fore_color.rgb
                                shape_info += f", fill_color=#{color[0]:02x}{color[1]:02x}{color[2]:02x}"
                            except:
                                pass
                except:
                    pass
                logger.info(shape_info)
            
            # Удаляем существующие текстовые поля и placeholder shapes
            shapes_to_remove = []
            for shape in slide.shapes:
                # Удаляем ВСЕ текстовые поля (TEXT_BOX) - они будут пересозданы
                # Это предотвращает дублирование текста при повторном сохранении
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    shapes_to_remove.append(shape)
                    text_preview = shape.text_frame.text[:30] if shape.has_text_frame and shape.text_frame.text else "empty"
                    logger.info(f"Marking TEXT_BOX for removal: {text_preview}")
                    continue  # Переходим к следующему shape
                
                # Удаляем shapes с placeholder'ами (если они не TEXT_BOX)
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if "{{" in text and "}}" in text:
                        shapes_to_remove.append(shape)
                        logger.info(f"Marking placeholder shape for removal: {text[:50]}")
                
                # Удаляем placeholder shapes (встроенные в PowerPoint)
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    # Не удаляем title и body placeholders если они нужны для layout
                    # Но удаляем content placeholders
                    try:
                        placeholder_type = shape.placeholder_format.type
                        # TYPE_OBJECT = 7, TYPE_CONTENT = 2
                        if placeholder_type in (2, 7):  
                            shapes_to_remove.append(shape)
                            logger.info(f"Marking content placeholder for removal: type={placeholder_type}")
                    except:
                        pass
                
                # Удаляем красные прямоугольники без текста (скорее всего, это артефакты)
                try:
                    if hasattr(shape, 'fill') and shape.shape_type == 1:  # AUTO_SHAPE
                        if shape.fill.type == 1:  # SOLID fill
                            try:
                                color = shape.fill.fore_color.rgb
                                # Проверяем, красный ли это цвет (R > 200, G < 100, B < 100)
                                if color[0] > 200 and color[1] < 100 and color[2] < 100:
                                    # Если нет текста или текст пустой - удаляем
                                    has_text = shape.has_text_frame and shape.text_frame.text.strip()
                                    if not has_text:
                                        shapes_to_remove.append(shape)
                                        logger.warning(f"Marking RED shape for removal: color=#{color[0]:02x}{color[1]:02x}{color[2]:02x}, name={shape.name if hasattr(shape, 'name') else 'unknown'}")
                            except:
                                pass
                except Exception as e:
                    logger.debug(f"Error checking shape fill: {e}")
            
            # Удаляем shapes
            for shape in shapes_to_remove:
                try:
                    sp = shape.element
                    sp.getparent().remove(sp)
                except Exception as e:
                    logger.warning(f"Failed to remove shape: {e}")
            
            logger.info(f"Removed {len(shapes_to_remove)} shapes (placeholders and text boxes)")
            
            # Фильтруем поля - убираем пустые значения
            filtered_fields = []
            for field in fields_data:
                text = field.get("text", "")
                # Пропускаем поля с пустым текстом или только пробелами
                if text and text.strip():
                    filtered_fields.append(field)
                else:
                    logger.debug(f"Skipping empty field at order_index {field.get('order_index', 0)}")
            
            logger.info(f"Filtered fields: {len(filtered_fields)} non-empty out of {len(fields_data)} total")
            
            # Сортируем поля по order_index
            sorted_fields = sorted(filtered_fields, key=lambda x: x.get("order_index", 0))
            
            # Используем TextLayoutService для размещения полей
            success = TextLayoutService.add_fields_with_auto_layout(
                slide, prs, sorted_fields
            )
            
            if success:
                logger.info(f"Successfully filled slide {slide_index} with {len(sorted_fields)} fields")
            
            return success
            
        except Exception as e:
            logger.error(f"Error filling template with data: {e}")
            return False
    
    @staticmethod
    def get_slide_count(prs: PPTXPresentation) -> int:
        """Get number of slides in presentation"""
        return len(prs.slides)
    
    @staticmethod
    def create_presentation_from_template(template_data: bytes) -> bytes:
        """
        Create a presentation from a master template
        
        Args:
            template_data: PPTX template file data as bytes
            
        Returns:
            PPTX file as bytes
        """
        # Load the template
        prs = PPTXService.load_presentation(template_data)
        
        logger.info(f"Created presentation from template with {len(prs.slides)} slide(s)")
        
        # Return the presentation
        return PPTXService.save_presentation(prs)
    
    @staticmethod
    def clone_slide(prs: PPTXPresentation, source_slide_index: int = 0) -> int:
        """
        Clone a slide within a presentation (preserves design and layout)
        
        Args:
            prs: Presentation object
            source_slide_index: Index of the slide to clone (default: 0 - first slide)
            
        Returns:
            Index of the newly created slide
        """
        if source_slide_index >= len(prs.slides):
            raise ValueError(f"Source slide index {source_slide_index} out of range")
        
        source_slide = prs.slides[source_slide_index]
        
        # Use the same layout as the source slide
        slide_layout = source_slide.slide_layout
        new_slide = prs.slides.add_slide(slide_layout)
        
        # Copy all shapes from the source slide
        for shape in source_slide.shapes:
            PPTXService._copy_shape(shape, new_slide)
        
        new_index = len(prs.slides) - 1
        logger.info(f"Cloned slide {source_slide_index} -> new slide at index {new_index}")
        return new_index
