"""
Text layout service for auto-sizing and positioning text fields
"""
import logging
from typing import Dict, List, Tuple, Optional
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)


class TextLayoutService:
    """Service for managing text layout with auto-sizing"""
    
    # Константы для отступов (в дюймах)
    MARGIN_LEFT = 1.0  # 1 дюйм слева
    MARGIN_RIGHT = 1.0  # 1 дюйм справа
    MARGIN_TOP = 1.0  # 1 дюйм сверху
    FIELD_VERTICAL_SPACING = 0.3  # Отступ между полями (в дюймах)
    LINE_SPACING_MULTIPLIER = 1.0  # Одинарный межстрочный интервал
    
    # Минимальный размер шрифта при уменьшении
    MIN_FONT_SIZE = 8
    
    @staticmethod
    def get_slide_dimensions(prs: PPTXPresentation) -> Tuple[float, float]:
        """
        Получить размеры слайда в дюймах
        
        Returns:
            (width, height) в дюймах
        """
        width_inches = prs.slide_width.inches
        height_inches = prs.slide_height.inches
        return width_inches, height_inches
    
    @staticmethod
    def calculate_text_box_width(slide_width: float) -> float:
        """
        Вычислить ширину текстового поля с учетом отступов
        
        Args:
            slide_width: Ширина слайда в дюймах
            
        Returns:
            Ширина текстового поля в дюймах
        """
        return slide_width - TextLayoutService.MARGIN_LEFT - TextLayoutService.MARGIN_RIGHT
    
    @staticmethod
    def estimate_text_height(
        text: str,
        font_size: float,
        box_width: float,
        font_name: str = "Arial"
    ) -> float:
        """
        Оценить высоту текстового поля для заданного текста
        
        Args:
            text: Текст для размещения
            font_size: Размер шрифта в пунктах
            box_width: Ширина текстового поля в дюймах
            font_name: Название шрифта
            
        Returns:
            Примерная высота в дюймах
        """
        # Примерная формула для расчета высоты
        # Используем среднюю ширину символа относительно размера шрифта
        
        # Средняя ширина символа (в дюймах) для данного размера шрифта
        avg_char_width = (font_size / 72) * 0.6  # 72 points per inch, 0.6 - коэффициент
        
        # Количество символов, которые помещаются в строку
        chars_per_line = int(box_width / avg_char_width)
        
        if chars_per_line <= 0:
            chars_per_line = 1
        
        # Количество строк
        num_lines = max(1, len(text) / chars_per_line)
        
        # Высота одной строки (размер шрифта + межстрочный интервал)
        line_height = (font_size / 72) * TextLayoutService.LINE_SPACING_MULTIPLIER * 1.2
        
        # Общая высота
        total_height = num_lines * line_height
        
        return total_height
    
    @staticmethod
    def create_text_field(
        slide,
        text: str,
        left: float,
        top: float,
        width: float,
        font_metadata: Dict,
        auto_height: bool = True
    ):
        """
        Создать текстовое поле с заданными параметрами
        
        Args:
            slide: Слайд, на который добавляется поле
            text: Текст для размещения
            left: Позиция слева в дюймах
            top: Позиция сверху в дюймах
            width: Ширина поля в дюймах
            font_metadata: Метаданные шрифта (name, size, color, bold, italic, alignment)
            auto_height: Автоматически рассчитать высоту
            
        Returns:
            Созданный text box shape
        """
        # Извлекаем метаданные шрифта
        font_name = font_metadata.get("name", "Arial") or "Arial"
        font_size = font_metadata.get("size", 18) or 18  # Защита от None
        font_color = font_metadata.get("color", "#000000") or "#000000"
        bold = font_metadata.get("bold", False) or False
        italic = font_metadata.get("italic", False) or False
        alignment = font_metadata.get("alignment", "LEFT") or "LEFT"
        
        # Проверяем, что font_size это число
        if not isinstance(font_size, (int, float)):
            logger.warning(f"Invalid font_size: {font_size}, using default 18")
            font_size = 18
        
        # Рассчитываем высоту, если нужно
        if auto_height:
            height = TextLayoutService.estimate_text_height(
                text, font_size, width, font_name
            )
        else:
            height = 1.0  # Дюйм по умолчанию
        
        # Создаем текстовое поле
        textbox = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height)
        )
        
        # Настраиваем текст
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        
        # Настраиваем форматирование
        if text_frame.paragraphs:
            paragraph = text_frame.paragraphs[0]
            
            # Выравнивание
            if alignment.upper() == "CENTER":
                paragraph.alignment = PP_ALIGN.CENTER
            elif alignment.upper() == "RIGHT":
                paragraph.alignment = PP_ALIGN.RIGHT
            else:
                paragraph.alignment = PP_ALIGN.LEFT
            
            # Форматирование шрифта
            if paragraph.runs:
                run = paragraph.runs[0]
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.italic = italic
                
                # Цвет (если указан)
                if font_color and font_color.startswith("#"):
                    try:
                        # Конвертируем hex в RGB
                        rgb_int = int(font_color[1:], 16)
                        from pptx.util import RGBColor
                        run.font.color.rgb = RGBColor(
                            (rgb_int >> 16) & 0xFF,
                            (rgb_int >> 8) & 0xFF,
                            rgb_int & 0xFF
                        )
                    except:
                        pass
        
        logger.info(f"Created text field: {len(text)} chars, {height:.2f}\" height")
        return textbox
    
    @staticmethod
    def layout_fields_vertically(
        slide,
        fields_data: List[Dict],
        slide_width: float,
        slide_height: float
    ) -> Tuple[List, float]:
        """
        Разместить поля вертикально одно за другим
        
        Args:
            slide: Слайд для размещения
            fields_data: Список данных полей [{text, font_metadata}, ...]
            slide_width: Ширина слайда в дюймах
            slide_height: Высота слайда в дюймах
            
        Returns:
            (list_of_shapes, total_height)
        """
        shapes = []
        current_top = TextLayoutService.MARGIN_TOP
        box_width = TextLayoutService.calculate_text_box_width(slide_width)
        
        for i, field_data in enumerate(fields_data):
            text = field_data.get("text", "")
            font_metadata = field_data.get("font_metadata", {})
            
            # Логируем данные поля для диагностики
            logger.info(f"Processing field {i+1}/{len(fields_data)}: text_len={len(text)}, font_metadata={font_metadata}")
            
            # Создаем поле
            shape = TextLayoutService.create_text_field(
                slide=slide,
                text=text,
                left=TextLayoutService.MARGIN_LEFT,
                top=current_top,
                width=box_width,
                font_metadata=font_metadata,
                auto_height=True
            )
            
            shapes.append(shape)
            
            # Обновляем позицию для следующего поля
            # Добавляем высоту текущего поля + вертикальный отступ между полями
            field_height = shape.height.inches
            current_top += field_height + TextLayoutService.FIELD_VERTICAL_SPACING
        
        total_height = current_top
        return shapes, total_height
    
    @staticmethod
    def shrink_fonts_if_needed(
        shapes: List,
        slide_height: float,
        total_height: float
    ) -> bool:
        """
        Уменьшить шрифты всех полей, если содержимое не помещается на слайде
        
        Args:
            shapes: Список text box shapes
            slide_height: Высота слайда в дюймах
            total_height: Общая высота содержимого в дюймах
            
        Returns:
            True если шрифты были уменьшены
        """
        # Проверяем, помещается ли содержимое
        if total_height <= slide_height:
            logger.info("Content fits on slide, no font shrinking needed")
            return False
        
        logger.warning(f"Content overflow: {total_height:.2f}\" > {slide_height:.2f}\"")
        
        # Вычисляем коэффициент уменьшения
        shrink_factor = slide_height / total_height
        
        # Уменьшаем шрифты всех полей
        for shape in shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            new_size = int(run.font.size.pt * shrink_factor)
                            # Не уменьшаем меньше минимального размера
                            new_size = max(new_size, TextLayoutService.MIN_FONT_SIZE)
                            run.font.size = Pt(new_size)
                            logger.info(f"Shrunk font: {run.font.size.pt}pt -> {new_size}pt")
        
        logger.info(f"All fonts shrunk by factor {shrink_factor:.2f}")
        return True
    
    @staticmethod
    def add_fields_with_auto_layout(
        slide,
        prs: PPTXPresentation,
        fields_data: List[Dict]
    ) -> bool:
        """
        Добавить поля на слайд с автоматическим размещением и масштабированием
        
        Args:
            slide: Слайд для размещения
            prs: Презентация (для получения размеров)
            fields_data: Список данных полей
            
        Returns:
            True если успешно
        """
        try:
            # Получаем размеры слайда
            slide_width, slide_height = TextLayoutService.get_slide_dimensions(prs)
            
            logger.info(f"Adding {len(fields_data)} fields to slide ({slide_width}\" x {slide_height}\")")
            
            # Размещаем поля вертикально
            shapes, total_height = TextLayoutService.layout_fields_vertically(
                slide, fields_data, slide_width, slide_height
            )
            
            # Уменьшаем шрифты, если не помещается
            TextLayoutService.shrink_fonts_if_needed(shapes, slide_height, total_height)
            
            logger.info("Fields added successfully with auto-layout")
            return True
            
        except Exception as e:
            logger.error(f"Error adding fields with auto-layout: {e}")
            return False
