from pptx import Presentation
from PIL import Image, ImageDraw
from io import BytesIO
from typing import BinaryIO
import logging

logger = logging.getLogger(__name__)


class PreviewGenerator:
    """Генератор PNG превью для блоков презентации"""
    
    # Стандартные размеры слайда в EMUs (English Metric Units)
    # 1 EMU = 1/914400 дюйма
    SLIDE_WIDTH_EMU = 9144000  # 10 дюймов
    SLIDE_HEIGHT_EMU = 6858000  # 7.5 дюймов
    
    # Размеры превью в пикселях
    PREVIEW_WIDTH = 800
    PREVIEW_HEIGHT = 600
    
    async def generate_preview(self, file_data: BinaryIO, slide_number: int = 1) -> BytesIO:
        """
        Генерирует превью для первого слайда презентации
        
        Args:
            file_data: BinaryIO данные PPTX файла
            slide_number: Номер слайда для превью (по умолчанию 1)
            
        Returns:
            BytesIO с PNG изображением
        """
        try:
            # Читаем данные
            if hasattr(file_data, 'read'):
                data = file_data.read()
                if hasattr(file_data, 'seek'):
                    file_data.seek(0)
            else:
                data = file_data
            
            # Генерируем превью
            preview_bytes = self.generate_slide_preview(data, slide_number)
            return BytesIO(preview_bytes)
            
        except Exception as e:
            logger.error(f"Error generating preview: {e}")
            # Возвращаем заглушку
            return BytesIO(self._generate_placeholder_preview())
    
    @staticmethod
    def generate_slide_preview(file_data: bytes, slide_number: int) -> bytes:
        """
        Генерирует превью для целого слайда
        
        Args:
            file_data: Бинарные данные PPTX файла
            slide_number: Номер слайда (1-based)
            
        Returns:
            PNG изображение в виде bytes
        """
        try:
            prs = Presentation(BytesIO(file_data))
            
            if slide_number < 1 or slide_number > len(prs.slides):
                raise ValueError(f"Slide number {slide_number} out of range")
            
            slide = prs.slides[slide_number - 1]
            
            # Создаем белое изображение
            img = Image.new('RGB', (PreviewGenerator.PREVIEW_WIDTH, PreviewGenerator.PREVIEW_HEIGHT), 'white')
            draw = ImageDraw.Draw(img)
            
            # Рисуем рамки для каждого shape на слайде
            for shape in slide.shapes:
                PreviewGenerator._draw_shape_placeholder(draw, shape)
            
            # Конвертируем в bytes
            output = BytesIO()
            img.save(output, format='PNG')
            output.seek(0)
            
            return output.read()
            
        except Exception as e:
            logger.error(f"Error generating slide preview: {e}")
            raise
    
    @staticmethod
    def generate_block_preview(file_data: bytes, slide_number: int, block_position: dict, block_size: dict) -> bytes:
        """
        Генерирует превью для отдельного блока
        
        Args:
            file_data: Бинарные данные PPTX файла
            slide_number: Номер слайда (1-based)
            block_position: Позиция блока {"x": int, "y": int} в EMUs
            block_size: Размер блока {"width": int, "height": int} в EMUs
            
        Returns:
            PNG изображение в виде bytes
        """
        try:
            # Генерируем превью всего слайда
            slide_preview = PreviewGenerator.generate_slide_preview(file_data, slide_number)
            
            # Загружаем изображение
            img = Image.open(BytesIO(slide_preview))
            
            # Конвертируем координаты из EMU в пиксели
            x_px = int((block_position["x"] / PreviewGenerator.SLIDE_WIDTH_EMU) * PreviewGenerator.PREVIEW_WIDTH)
            y_px = int((block_position["y"] / PreviewGenerator.SLIDE_HEIGHT_EMU) * PreviewGenerator.PREVIEW_HEIGHT)
            w_px = int((block_size["width"] / PreviewGenerator.SLIDE_WIDTH_EMU) * PreviewGenerator.PREVIEW_WIDTH)
            h_px = int((block_size["height"] / PreviewGenerator.SLIDE_HEIGHT_EMU) * PreviewGenerator.PREVIEW_HEIGHT)
            
            # Обрезаем изображение по координатам блока
            # Учитываем границы изображения
            x_px = max(0, x_px)
            y_px = max(0, y_px)
            x2_px = min(PreviewGenerator.PREVIEW_WIDTH, x_px + w_px)
            y2_px = min(PreviewGenerator.PREVIEW_HEIGHT, y_px + h_px)
            
            if x2_px <= x_px or y2_px <= y_px:
                # Если координаты невалидны, создаем пустое изображение
                cropped = Image.new('RGB', (200, 150), 'lightgray')
            else:
                cropped = img.crop((x_px, y_px, x2_px, y2_px))
            
            # Конвертируем в bytes
            output = BytesIO()
            cropped.save(output, format='PNG')
            output.seek(0)
            
            return output.read()
            
        except Exception as e:
            logger.error(f"Error generating block preview: {e}")
            # Возвращаем заглушку в случае ошибки
            return PreviewGenerator._generate_placeholder_preview()
    
    @staticmethod
    def _draw_shape_placeholder(draw: ImageDraw.ImageDraw, shape):
        """Рисует заглушку для shape"""
        try:
            # Конвертируем координаты из EMU в пиксели
            x = int((shape.left / PreviewGenerator.SLIDE_WIDTH_EMU) * PreviewGenerator.PREVIEW_WIDTH)
            y = int((shape.top / PreviewGenerator.SLIDE_HEIGHT_EMU) * PreviewGenerator.PREVIEW_HEIGHT)
            w = int((shape.width / PreviewGenerator.SLIDE_WIDTH_EMU) * PreviewGenerator.PREVIEW_WIDTH)
            h = int((shape.height / PreviewGenerator.SLIDE_HEIGHT_EMU) * PreviewGenerator.PREVIEW_HEIGHT)
            
            # Рисуем прямоугольник
            draw.rectangle([x, y, x + w, y + h], outline='gray', width=1)
            
            # Если есть текст, пытаемся его отобразить
            if hasattr(shape, 'text') and shape.text:
                # Упрощенное отображение текста (для полноценного нужен PIL с поддержкой шрифтов)
                text = shape.text[:50]  # Ограничиваем длину
                draw.text((x + 5, y + 5), text, fill='black')
                
        except Exception as e:
            logger.warning(f"Error drawing shape placeholder: {e}")
    
    @staticmethod
    def _generate_placeholder_preview() -> bytes:
        """Генерирует placeholder изображение"""
        img = Image.new('RGB', (200, 150), 'lightgray')
        draw = ImageDraw.Draw(img)
        draw.text((50, 70), "Preview N/A", fill='darkgray')
        
        output = BytesIO()
        img.save(output, format='PNG')
        output.seek(0)
        return output.read()
