from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Dict, Any, Tuple, Optional
import logging
import re
from io import BytesIO

logger = logging.getLogger(__name__)


class PPTXParser:
    """Парсер PPTX файлов для извлечения блоков"""
    
    @staticmethod
    def _extract_placeholders_from_text(text: str) -> List[Dict[str, str]]:
        """
        Извлекает плейсхолдеры из текста в формате {{label.name}}
        
        Args:
            text: Текст для анализа
            
        Returns:
            Список словарей с ключами: 'placeholder' (полный текст), 'label', 'name'
            
        Примеры:
            {{title.about_project}} -> label='title', name='about_project', placeholder='{{title.about_project}}'
            {{subtitle.project_description}} -> label='subtitle', name='project_description'
            {{list.project_goals}} -> label='list', name='project_goals'
            {{text.project_description}} -> label='text', name='project_description'
        """
        pattern = r'\{\{([^.}]+)\.([^}]+)\}\}'
        placeholders = []
        
        for match in re.finditer(pattern, text):
            placeholders.append({
                'placeholder': match.group(0),  # {{title.about_project}}
                'label': match.group(1).strip(),  # title (тип поля)
                'name': match.group(2).strip()  # about_project (имя поля)
            })
        
        return placeholders
    
    @staticmethod
    def parse_presentation(file_data: bytes) -> List[Dict[str, Any]]:
        """
        Парсит PPTX файл и извлекает блоки
        
        Args:
            file_data: Бинарные данные PPTX файла
            
        Returns:
            Список блоков с метаданными
        """
        try:
            prs = Presentation(BytesIO(file_data))
            blocks = []
            
            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_blocks = PPTXParser._parse_slide(slide, slide_idx)
                blocks.extend(slide_blocks)
            
            logger.info(f"Parsed {len(blocks)} blocks from {len(prs.slides)} slides")
            return blocks
            
        except Exception as e:
            logger.error(f"Error parsing PPTX: {e}")
            raise
    
    @staticmethod
    def _parse_slide(slide, slide_number: int) -> List[Dict[str, Any]]:
        """Парсит один слайд и извлекает блоки"""
        blocks = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                block = PPTXParser._parse_shape(shape, slide_number, shape_idx)
                if block:
                    blocks.append(block)
            except Exception as e:
                logger.warning(f"Error parsing shape {shape_idx} on slide {slide_number}: {e}")
                continue
        
        return blocks
    
    @staticmethod
    def _parse_shape(shape, slide_number: int, shape_idx: int) -> Dict[str, Any]:
        """Парсит отдельный shape и извлекает информацию"""
        
        # Определяем тип блока
        block_type = PPTXParser._determine_block_type(shape)
        
        # Генерируем ключ блока
        key = PPTXParser._generate_block_key(shape, block_type, slide_number, shape_idx)
        
        # Извлекаем позицию и размер
        position = {
            "x": shape.left,
            "y": shape.top
        }
        
        size = {
            "width": shape.width,
            "height": shape.height
        }
        
        # Извлекаем метаданные
        block_metadata = PPTXParser._extract_metadata(shape, block_type)
        
        # Формируем схему данных
        value_schema = PPTXParser._generate_value_schema(shape, block_type)
        
        # Извлекаем плейсхолдеры из текста фигуры
        placeholders = []
        if hasattr(shape, "text_frame") and shape.has_text_frame:
            text = shape.text_frame.text
            placeholders = PPTXParser._extract_placeholders_from_text(text)
        
        return {
            "slide_number": slide_number,
            "type": block_type,
            "key": key,
            "value_schema": value_schema,
            "position": position,
            "size": size,
            "block_metadata": block_metadata,
            "placeholders": placeholders
        }
    
    @staticmethod
    def _determine_block_type(shape) -> str:
        """Определяет тип блока на основе shape"""
        
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "table"
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return "chart"
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            return "placeholder"
        elif hasattr(shape, "text_frame") and shape.has_text_frame:
            # Проверяем, является ли это списком
            if shape.text_frame.text and any(
                para.text.strip().startswith(('•', '-', '*', '1.', '2.'))
                for para in shape.text_frame.paragraphs
            ):
                return "list"
            return "text"
        elif shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
            return "shape"
        else:
            return "shape"
    
    @staticmethod
    def _generate_block_key(shape, block_type: str, slide_number: int, shape_idx: int) -> str:
        """Генерирует уникальный ключ для блока"""
        
        # Пытаемся использовать имя placeholder'а
        if hasattr(shape, 'is_placeholder'):
            try:
                if shape.is_placeholder:
                    placeholder_type = shape.placeholder_format.type
                    return f"{block_type}_{placeholder_type}_{slide_number}"
            except:
                pass
        
        # Пытаемся использовать имя shape
        if hasattr(shape, 'name') and shape.name:
            safe_name = shape.name.replace(' ', '_').lower()
            return f"{block_type}_{safe_name}"
        
        # Используем порядковый номер
        return f"{block_type}_s{slide_number}_i{shape_idx}"
    
    @staticmethod
    def _extract_metadata(shape, block_type: str) -> Dict[str, Any]:
        """
        Извлекает метаданные из shape
        
        Включает:
        - Информацию о шрифте (name, size, color, bold, italic, alignment)
        - Позицию и размеры shape
        - Placeholder информацию
        """
        metadata = {
            "shape_type": str(shape.shape_type),
            "shape_name": shape.name if hasattr(shape, 'name') else None,
        }
        
        # Сохраняем позицию и размеры (в EMU - English Metric Units)
        metadata["position"] = {
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height
        }
        
        # Для текстовых блоков извлекаем информацию о шрифте
        if block_type in ("text", "list", "placeholder") and hasattr(shape, "text_frame"):
            try:
                if shape.text_frame.paragraphs:
                    first_para = shape.text_frame.paragraphs[0]
                    
                    # Извлекаем alignment параграфа
                    metadata["alignment"] = str(first_para.alignment) if first_para.alignment else "LEFT"
                    
                    if first_para.runs:
                        first_run = first_para.runs[0]
                        font = first_run.font
                        
                        # Извлекаем цвет шрифта
                        font_color = None
                        try:
                            if font.color.type == 1:  # RGB
                                rgb = font.color.rgb
                                font_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                        except:
                            pass
                        
                        metadata["font"] = {
                            "name": font.name,
                            "size": font.size.pt if font.size else None,
                            "color": font_color,
                            "bold": font.bold,
                            "italic": font.italic,
                        }
            except Exception as e:
                logger.warning(f"Error extracting font metadata: {e}")
        
        # Для placeholder'ов
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            try:
                metadata["placeholder"] = {
                    "type": str(shape.placeholder_format.type),
                    "idx": shape.placeholder_format.idx
                }
            except:
                pass
        
        return metadata
    
    @staticmethod
    def _generate_value_schema(shape, block_type: str) -> Dict[str, Any]:
        """Генерирует схему данных для блока"""
        
        schema = {
            "type": block_type
        }
        
        if block_type in ("text", "placeholder"):
            schema["structure"] = {
                "text": "string"
            }
        elif block_type == "list":
            schema["structure"] = {
                "items": ["string"]
            }
        elif block_type == "image":
            schema["structure"] = {
                "image_url": "string",
                "alt_text": "string"
            }
        elif block_type == "table":
            try:
                if hasattr(shape, 'table'):
                    table = shape.table
                    schema["structure"] = {
                        "rows": table.rows.__len__() if hasattr(table, 'rows') else 0,
                        "columns": table.columns.__len__() if hasattr(table, 'columns') else 0,
                        "data": "array<array<string>>"
                    }
            except:
                schema["structure"] = {"data": "array<array<string>>"}
        elif block_type == "chart":
            schema["structure"] = {
                "chart_type": "string",
                "data": "object"
            }
        else:
            schema["structure"] = {}
        
        return schema
