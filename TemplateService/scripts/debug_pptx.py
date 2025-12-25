#!/usr/bin/env python3
"""
Отладочный скрипт для анализа PPTX файла
"""

import sys
from pptx import Presentation

def analyze_pptx(file_path: str):
    """Анализирует структуру PPTX файла"""
    print(f"📄 Анализ файла: {file_path}")
    print("=" * 70)
    
    try:
        prs = Presentation(file_path)
        print(f"✅ Файл успешно открыт")
        print(f"📊 Количество слайдов: {len(prs.slides)}")
        print()
        
        total_shapes = 0
        for slide_idx, slide in enumerate(prs.slides, start=1):
            print(f"🔹 Слайд {slide_idx}:")
            print(f"   Количество shapes: {len(slide.shapes)}")
            
            if len(slide.shapes) == 0:
                print(f"   ⚠️  ПУСТОЙ СЛАЙД!")
                continue
            
            for shape_idx, shape in enumerate(slide.shapes):
                total_shapes += 1
                print(f"   Shape {shape_idx + 1}:")
                print(f"      Name: {shape.name if hasattr(shape, 'name') else 'N/A'}")
                print(f"      Type: {shape.shape_type}")
                print(f"      Has text: {hasattr(shape, 'text')}")
                
                if hasattr(shape, 'text'):
                    text = shape.text[:50] if shape.text else "(пусто)"
                    print(f"      Text: {text}")
                
                if hasattr(shape, 'shape_type'):
                    print(f"      Shape type value: {shape.shape_type}")
                
                # Проверяем текстовый фрейм
                if hasattr(shape, 'has_text_frame'):
                    print(f"      Has text frame: {shape.has_text_frame}")
                
                # Проверяем placeholder
                if hasattr(shape, 'is_placeholder'):
                    print(f"      Is placeholder: {shape.is_placeholder}")
                    if shape.is_placeholder:
                        try:
                            print(f"      Placeholder type: {shape.placeholder_format.type}")
                        except:
                            pass
                
                print(f"      Position: left={shape.left}, top={shape.top}")
                print(f"      Size: width={shape.width}, height={shape.height}")
                print()
        
        print("=" * 70)
        print(f"📊 ИТОГО:")
        print(f"   Слайдов: {len(prs.slides)}")
        print(f"   Shapes: {total_shapes}")
        print()
        
        if total_shapes == 0:
            print("❌ ПРОБЛЕМА: В файле нет shapes!")
            print()
            print("💡 Возможные причины:")
            print("   1. Слайды пустые")
            print("   2. Все элементы сгруппированы")
            print("   3. Файл повреждён")
            print()
            print("🔧 Решение:")
            print("   1. Откройте PPTX в PowerPoint")
            print("   2. Добавьте текстовые блоки, изображения или фигуры")
            print("   3. Убедитесь, что элементы не сгруппированы")
            print("   4. Сохраните файл")
        
    except Exception as e:
        print(f"❌ Ошибка: {e}")
        import traceback
        traceback.print_exc()


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Использование: python debug_pptx.py <путь_к_файлу.pptx>")
        sys.exit(1)
    
    analyze_pptx(sys.argv[1])
